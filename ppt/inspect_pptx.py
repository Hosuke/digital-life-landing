"""Inspect the structure of a PPTX file."""
from pptx import Presentation
import sys

pptx_file = sys.argv[1] if len(sys.argv) > 1 else "总决赛ppt.pptx"
prs = Presentation(pptx_file)

print(f"Slide size: {prs.slide_width}x{prs.slide_height} EMU")
print(f"  = {prs.slide_width/914400:.1f} x {prs.slide_height/914400:.1f} inches")
print(f"Slides: {len(prs.slides)}")
print(f"Layouts: {len(prs.slide_layouts)}")
print()

for idx, slide in enumerate(prs.slides, 1):
    print(f"=== Slide {idx} (layout: {slide.slide_layout.name}) ===")
    for shape in slide.shapes:
        has_text = shape.has_text_frame
        txt_preview = ""
        if has_text:
            txt_preview = shape.text_frame.text.replace("\n", " | ")[:100]
        print(f"  [{shape.shape_type}] \"{shape.name}\"  text=\"{txt_preview}\"")
        if has_text:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                p_text = p.text.strip()
                if not p_text:
                    continue
                for r in p.runs:
                    fn = r.font.name
                    fs = r.font.size
                    fb = r.font.bold
                    try:
                        fc = str(r.font.color.rgb) if r.font.color and r.font.color.rgb else None
                    except:
                        fc = None
                    print(f"    p{pi} run: \"{r.text[:50]}\" font={fn} size={fs} bold={fb} color={fc}")
    print()
