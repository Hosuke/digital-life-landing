"""Extract text outline from a PPTX file into Markdown format."""

import sys
from pptx import Presentation


def extract_outline(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    lines = []

    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {idx}")

        # Try to get slide layout name
        try:
            layout_name = slide.slide_layout.name
            lines.append(f"_Layout: {layout_name}_\n")
        except Exception:
            lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Rough heuristic: larger or bold = heading
                    is_bold = any(
                        run.font.bold for run in para.runs if run.font.bold is not None
                    )
                    level = para.level if para.level else 0
                    indent = "  " * level
                    if is_bold:
                        lines.append(f"{indent}**{text}**")
                    else:
                        lines.append(f"{indent}- {text}")

            # Note tables
            if shape.has_table:
                table = shape.table
                lines.append(f"\n_[Table: {table.rows.__len__()} rows x {len(table.columns)} cols]_")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")

        lines.append("")  # blank line between slides
        lines.append("---")
        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "总决赛ppt.pptx"
    outline = extract_outline(pptx_file)
    out_path = pptx_file.rsplit(".", 1)[0] + "_outline.md"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(outline)
    print(f"Outline saved to {out_path}")
