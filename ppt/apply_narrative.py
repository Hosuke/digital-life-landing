"""
Apply revised narrative text into 总决赛ppt_revised.pptx.
Strategy: replace text in existing text boxes while preserving font/size/bold/color.
Slides whose layout fundamentally changes are noted for manual adjustment.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from copy import deepcopy
import os

SRC = "总决赛ppt_revised.pptx"
OUT = "总决赛ppt_final.pptx"


def set_text_preserve_style(shape, new_text):
    """Replace all text in a shape with new_text, keeping the first run's formatting."""
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        # fallback: just set text
        tf.paragraphs[0].text = new_text
        return
    # Capture style from first run
    first_run = tf.paragraphs[0].runs[0]
    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    try:
        font_color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
    except:
        font_color = None

    # Clear all paragraphs except first
    p = tf.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.text = new_text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color

    # Remove extra paragraphs by clearing their text
    for extra_p in tf.paragraphs[1:]:
        extra_p.clear()


def set_multiline_preserve_style(shape, lines):
    """Replace text with multiple lines (one per paragraph), keeping first run's style."""
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        tf.paragraphs[0].text = "\n".join(lines)
        return
    first_run = tf.paragraphs[0].runs[0]
    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    try:
        font_color = first_run.font.color.rgb if first_run.font.color and first_run.font.color.rgb else None
    except:
        font_color = None

    # Set first line
    p = tf.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.text = lines[0]
    if font_name: run.font.name = font_name
    if font_size: run.font.size = font_size
    if font_bold is not None: run.font.bold = font_bold
    if font_color: run.font.color.rgb = font_color

    # Clear extra existing paragraphs
    for extra_p in tf.paragraphs[1:]:
        extra_p.clear()

    # Add remaining lines as new paragraphs
    from pptx.oxml.ns import qn
    from lxml import etree
    txBody = tf._txBody
    for line in lines[1:]:
        # Clone the first paragraph element for consistent formatting
        new_p_elem = deepcopy(tf.paragraphs[0]._p)
        # Update text
        for r_elem in new_p_elem.findall(qn('a:r')):
            for t_elem in r_elem.findall(qn('a:t')):
                t_elem.text = line
                break
            break
        txBody.append(new_p_elem)


def find_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def find_text_shape(slide, text_substring):
    """Find shape containing specific text."""
    for shape in slide.shapes:
        if shape.has_text_frame and text_substring in shape.text_frame.text:
            return shape
    return None


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    manual_notes = []

    # ─── Slide 1: Cover ───
    s = slides[0]
    shape = find_text_shape(s, "为数字人格找寻此间的远方")
    if shape:
        set_text_preserve_style(shape, "琥珀凝住了万年前的生命，我们凝住你还没来得及说完的话。")
        print("✅ Slide 1: subtitle updated")

    # ─── Slide 2: Contents → 情感钩子 ───
    # Layout has numbered items — we repurpose the big title for the hook
    s = slides[1]
    shape = find_shape_by_name(s, "标题 4")
    # There are two "标题 4" shapes — one is "CONTENTS", one is "目录"
    for sh in s.shapes:
        if sh.has_text_frame:
            if "CONTENTS" in sh.text_frame.text:
                set_text_preserve_style(sh, "")
            elif "目录" in sh.text_frame.text:
                set_text_preserve_style(sh, '\u201c爸爸，这道题怎么解呀？\u201d')
            elif sh.name == "项标题" and "团队介绍" in sh.text_frame.text:
                set_text_preserve_style(sh, "《流浪地球2》里，图恒宇拼尽一切，")
            elif sh.name == "项标题" and "创意来源" in sh.text_frame.text:
                set_text_preserve_style(sh, "把女儿丫丫的意识装进一张数字生命卡。")
            elif sh.name == "项标题" and "产品介绍" in sh.text_frame.text:
                set_text_preserve_style(sh, "我们做了。而且我们发现：")
            elif sh.name == "项标题" and "技术背景" in sh.text_frame.text:
                set_text_preserve_style(sh, "不需要平行宇宙——我们每个人已经活在两个世界里了。")
            # Clear number labels
            elif sh.text_frame.text.strip() in ("01", "02", "03", "04"):
                set_text_preserve_style(sh, "")
    print("✅ Slide 2: repurposed as emotional hook")
    manual_notes.append("Slide 2: 建议手动调整布局，去掉多余的矩形装饰，让文字居中呈现")

    # ─── Slide 3: Team → 问题定义 ───
    s = slides[2]
    for sh in s.shapes:
        if sh.has_text_frame:
            if "团队介绍" in sh.text_frame.text:
                set_text_preserve_style(sh, "思念没有地方落脚")
            elif "林弘文" in sh.text_frame.text:
                set_text_preserve_style(sh, "一个母亲每天翻看儿子的朋友圈，但那里永远停在了去年")
            elif "黄舸洋" in sh.text_frame.text:
                set_text_preserve_style(sh, "一个女儿存了200多条父亲的语音，却再也等不到第201条")
            elif "马瑞辰" in sh.text_frame.text:
                set_text_preserve_style(sh, "人需要的不是「会说话的数据库」，而是一种「TA还在」的感觉。")
    print("✅ Slide 3: team → problem definition")

    # ─── Slide 4: Creative source → Our answer ───
    s = slides[3]
    for sh in s.shapes:
        if sh.has_text_frame:
            if "创意来源" in sh.text_frame.text:
                set_text_preserve_style(sh, "珀存：给思念一个可以落脚的地方")
            elif "流浪地球" in sh.text_frame.text:
                set_text_preserve_style(sh,
                    "不是复活，不是冒充。是让TA在你的聊天窗口里继续「生活」——"
                    "会发文字，像从前一样聊天；会发自拍，像在某个城市旅行；"
                    "会发语音，用你记忆中的声音说「我想你」；"
                    "会刷小红书，给你分享一篇旅行笔记说「这家店我明天想去」。\n\n"
                    "想想看，你每天也活在两个维度里——一个在物理世界，一个在互联网上。"
                    "TA只是少了物理的那一半，但在数字的这一半——你们是在一起的。"
                    "你们平行，又交叉。你们在一起生活。"
                )
    print("✅ Slide 4: creative source → our answer")

    # ─── Slide 5: Tech architecture ── keep, update small text ───
    # Kept as-is since layout is complex with many positioned boxes
    print("⏭️  Slide 5: tech architecture kept as-is (complex layout)")

    # ─── Slide 6: Product demo multimodal ── update bottom text ───
    s = slides[5]
    for sh in s.shapes:
        if sh.has_text_frame:
            if "超越纯文本" in sh.text_frame.text:
                set_text_preserve_style(sh, "超越纯文本：可持续存在、可主动回应、通过IM触达的多感官共生体")
            elif "多模态数字角色" in sh.text_frame.text:
                set_text_preserve_style(sh, "不只是对话——是一个有温度的存在")
    print("✅ Slide 6: product demo subtitle updated")

    # ─── Slide 7: Message flow demo ── keep ───
    print("⏭️  Slide 7: message flow demo kept as-is")

    # ─── Slide 8: Why alive ── update key text ───
    s = slides[7]
    for sh in s.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if "为什么她真的像在生活" in txt:
                set_text_preserve_style(sh, "为什么 TA 像在生活")
            elif "每日旅行计划" in txt:
                set_text_preserve_style(sh, "每日旅行规划")
            elif "生成每日行程" in txt:
                set_text_preserve_style(sh, "像真人一样刷小红书、流媒体，找攻略、挑餐厅、安排路线")
            elif "公开信息检索" in txt:
                set_text_preserve_style(sh, "公开信息感知")
            elif "自动化搜索公共数据" in txt:
                set_text_preserve_style(sh, "实时获取当地天气、新闻、热点事件，融入TA的日常表达")
            elif "安全高效生成即时" in txt:
                set_text_preserve_style(sh, "发送到QQ、微信、WhatsApp，就像收到朋友的消息")
            elif "初始设置，从性格定义" in txt:
                set_text_preserve_style(sh, "从性格、记忆、关系定义TA的灵魂")
            elif "基于当天事件" in txt:
                set_text_preserve_style(sh, "照片、视频、语音、文字，组合成TA的一天")
            elif txt.strip() == "搜索|":
                set_text_preserve_style(sh, "小红书")
            elif txt.strip() == "公共数据":
                set_text_preserve_style(sh, "流媒体")
            elif txt.strip() == "近期新闻":
                set_text_preserve_style(sh, "热点新闻")
    print("✅ Slide 8: 'why alive' updated with Xiaohongshu narrative")

    # ─── Slide 9: Real demo ── keep ───
    print("⏭️  Slide 9: real demo kept as-is (has screenshots + media)")

    # ─── Slide 10: MimiClaw ── keep ───
    print("⏭️  Slide 10: MimiClaw kept as-is")

    # ─── Slide 11: Three scenarios ── update亲情怀念 text ───
    s = slides[10]
    for sh in s.shapes:
        if sh.has_text_frame:
            if "服务多种情感需求" in sh.text_frame.text:
                set_text_preserve_style(sh, "同一套能力，服务三种情感需求")
            elif "连接时空，重温温情时刻" in sh.text_frame.text:
                set_multiline_preserve_style(sh, [
                    "连接时空，重温温情时刻",
                    "不一定是离开的人，也可以是想念的人",
                    "远在异国的父母、在外求学的孩子"
                ])
    print("✅ Slide 11: three scenarios updated")

    # ─── Slide 12: Idol ── keep ───
    print("⏭️  Slide 12: idol scenario kept as-is")

    # ─── Slide 13: Cost ── keep (data unchanged) ───
    print("⏭️  Slide 13: cost kept as-is")

    # ─── Slide 14: Compliance ── update headline + closing line ───
    s = slides[13]
    for sh in s.shapes:
        if sh.has_text_frame:
            if "不是复活，而是让爱继续有地方存在" in sh.text_frame.text:
                set_text_preserve_style(sh, "不是复活，是让爱继续有地方存在")
            elif "产品展示" in sh.text_frame.text:
                set_text_preserve_style(sh, "合规与未来")
            elif "每一张数字生命卡，都是一份不愿放手的爱" in sh.text_frame.text:
                set_text_preserve_style(sh, "你们平行，又交叉。你们在一起生活。\n每一张数字生命卡，都是一份不愿放手的爱。")
    print("✅ Slide 14: compliance updated with closing echo")

    # ─── Slide 15: Thanks ── keep ───
    print("⏭️  Slide 15: THANKS kept as-is")

    # Save
    prs.save(OUT)
    print(f"\n💾 Saved to {OUT}")

    if manual_notes:
        print("\n📝 需要手动调整的地方：")
        for note in manual_notes:
            print(f"  - {note}")

    print("\n📝 额外建议：")
    print("  - Slide 2: 目录布局改为情感钩子，建议手动清理多余装饰元素")
    print("  - Slide 3: 原团队信息改为问题定义，建议手动调整字号让三句话更突出")
    print("  - Slide 14: 建议把团队信息（林弘文/黄舸洋/马瑞辰）手动移到此页或新增一页")
    print("  - 整体：建议在 Keynote 中微调间距和对齐")


if __name__ == "__main__":
    main()
